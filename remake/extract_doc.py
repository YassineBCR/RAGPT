import PyPDF2

def extract_text_from_pdf(file_path):
    pdf_file_obj = open(file_path, 'rb')
    pdf_reader = PyPDF2.PdfReader(pdf_file_obj) # utilisez PdfReader au lieu de PdfFileReader
    text = ''
    for page_obj in pdf_reader.pages: # itérez directement sur les pages au lieu d'utiliser numPages
        text += page_obj.extract_text()
    pdf_file_obj.close()
    return text
import openpyxl

def extract_text_from_xlsx(file_path):
    workbook = openpyxl.load_workbook(file_path)
    text = ''
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows():
            for cell in row:
                text += ' {} ' .format(cell.value)
    return text
#doc
import docx

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = ''.join([paragraph.text for paragraph in doc.paragraphs])
    return text
#powerpoint
from pptx import Presentation

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text
    return text

#txt
def extract_text_from_txt(file_path):
    with open(file_path, 'r') as file:
        return file.read()
import os 
def main():
    # Demander à l'utilisateur de choisir un fichier
    file_path = input("Entrez le chemin d'accès du fichier que vous voulez analyser : ")

    # Vérifier si le fichier existe
    if not os.path.isfile(file_path):
        print("Le fichier n'existe pas.")
        return

    # Déterminer le type de fichier en fonction de l'extension
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext == ".pdf":
        extracted_text = extract_text_from_pdf(file_path)
    elif file_ext == ".docx":
        extracted_text = extract_text_from_docx(file_path)
    elif file_ext == ".xlsx" or file_ext == ".xls":
        extracted_text = extract_text_from_xlsx(file_path)
    elif file_ext == ".pptx" or file_ext == ".ppt":
        extracted_text = extract_text_from_pptx(file_path)
    elif file_ext == ".txt":
        extracted_text = extract_text_from_txt(file_path)
    else:
        print("Le type de fichier n'est pas pris en charge.")
        return

    # Afficher le texte extrait
    print("Le contenu du fichier est : ")
    print(extracted_text)

if __name__ == "__main__":
    main()    